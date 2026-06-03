#!/usr/bin/env python3
"""Generate AgentRun Hackathon Pitch Deck PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ─────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_CARD     = RGBColor(0x16, 0x2A, 0x40)   # card background
ACCENT_BLUE = RGBColor(0x00, 0x8C, 0xFF)   # RC blue
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xBE)   # teal highlight
TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GREY   = RGBColor(0xA8, 0xBB, 0xCC)
TEXT_DIM    = RGBColor(0x5A, 0x7A, 0x94)
GREEN_OK    = RGBColor(0x00, 0xC8, 0x6E)
RED_WARN    = RGBColor(0xFF, 0x4D, 0x4D)
YELLOW      = RGBColor(0xFF, 0xC1, 0x07)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helper functions ─────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=BG_DARK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, radius=False):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=TEXT_WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def label(slide, text, x, y, w, h, bg_color=ACCENT_BLUE, fg_color=TEXT_WHITE, size=12):
    """Pill-style label."""
    r = rect(slide, x, y, w, h, bg_color)
    txt(slide, text, x + 0.05, y + 0.02, w - 0.1, h - 0.04,
        size=size, bold=True, color=fg_color, align=PP_ALIGN.CENTER)
    return r

def divider(slide, y, color=ACCENT_BLUE, alpha_width=12.0):
    """Horizontal accent line."""
    shape = slide.shapes.add_shape(1, Inches(0.66), Inches(y), Inches(alpha_width), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

# ── Slide 1: Title ───────────────────────────────────────────────

s = add_slide()
bg(s)

# Left accent bar
rect(s, 0, 0, 0.12, 7.5, ACCENT_BLUE)

# Tagline pill
label(s, "  Hackathon 2026  ", 0.5, 0.5, 2.2, 0.38, ACCENT_TEAL, BG_DARK, 11)

# Main title
txt(s, "AgentRun", 0.5, 1.1, 9, 1.6, size=80, bold=True, color=TEXT_WHITE)

# Subtitle
txt(s, "让 Agent 成为你 RC 团队里真正的同事", 0.5, 2.75, 9, 0.6,
    size=26, color=ACCENT_TEAL, bold=True)

txt(s, "不只是对话，而是真正执行", 0.5, 3.32, 9, 0.5,
    size=20, color=TEXT_GREY)

divider(s, 4.2)

# Bottom row
txt(s, "Built on RingCentral Platform", 0.5, 4.4, 5, 0.4, size=14, color=TEXT_DIM)
txt(s, "SMS · Fax · Phone · Task · Video", 0.5, 4.78, 6, 0.4, size=13, color=TEXT_DIM)

# Right decoration — three capability dots
for i, (label_text, col) in enumerate([
    ("Role Bot", ACCENT_BLUE),
    ("Human→Agent", ACCENT_TEAL),
    ("Agent→Agent", GREEN_OK),
]):
    xi = 10.5 + i * 0.9
    rect(s, xi, 2.8, 0.7, 0.7, col)
    txt(s, label_text, xi - 0.1, 3.6, 1.0, 0.4, size=9, color=TEXT_GREY, align=PP_ALIGN.CENTER)

add_notes(s, "开场：其他 AI 帮你想，我们的 Agent 真的执行——SMS 真发出去，传真真送到，电话真打出来。")

# ── Slide 2: Problem ─────────────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, RED_WARN)

txt(s, "问题", 0.5, 0.3, 3, 0.5, size=13, color=RED_WARN, bold=True)
txt(s, "AI 卡在了个人层面", 0.5, 0.65, 10, 0.9, size=42, bold=True)

# Flow diagram
boxes = [
    ("每个人在用 AI", 0.5),
    ("各自的 AI 对话", 3.0),
    ("组织协作靠人", 5.5),
    ("上下文丢失", 8.0),
]
for label_t, xi in boxes:
    rect(s, xi, 2.0, 2.2, 0.8, BG_CARD)
    txt(s, label_t, xi + 0.1, 2.1, 2.0, 0.6, size=14, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# Arrows between boxes
for xi in [2.7, 5.2, 7.7]:
    txt(s, "→", xi, 2.1, 0.5, 0.6, size=20, color=TEXT_DIM, align=PP_ALIGN.CENTER)

divider(s, 3.3, RED_WARN, 12.0)

# Four gaps
gaps = [
    ("❌  AI 如何真正发 SMS 给客户", 0.5),
    ("❌  AI 如何真正传真给 Lowe's HQ", 0.5),
    ("❌  AI 如何真的帮你打出电话", 6.5),
    ("❌  各部门 AI 如何自动协作", 6.5),
]
for i, (g, xi) in enumerate(gaps):
    yi = 3.7 + (i % 2) * 0.75
    txt(s, g, xi, yi, 5.5, 0.55, size=16, color=TEXT_WHITE)

txt(s, "Copilot  ·  Gemini  ·  ChatGPT  解决个人生产力，没有人解决执行层",
    0.5, 6.4, 12.0, 0.5, size=14, color=TEXT_DIM, italic=True)

add_notes(s, "ChatGPT 帮你起草邮件，但它不能帮你发。Copilot 帮你总结会议，但它不知道出了事该找谁。我们做的是执行层。")

# ── Slide 3: Real Customer ───────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_TEAL)

txt(s, "真实场景", 0.5, 0.3, 3, 0.5, size=13, color=ACCENT_TEAL, bold=True)
txt(s, "Keller Interiors", 0.5, 0.65, 10, 0.9, size=42, bold=True)

# Stats bar
stats = [("33 门店", "15 州", None), ("100-399 员工", None, None), ("Lowe's 27 年合作", None, None)]
for i, (a, b, c) in enumerate([
    ("33 门店", "美国 15 个州", ""),
    ("100-399 人", "地板安装服务商", ""),
    ("Lowe's 合作", "27 年，主要收入", ""),
]):
    xi = 0.5 + i * 4.2
    rect(s, xi, 1.85, 3.8, 1.0, BG_CARD)
    txt(s, a, xi + 0.15, 1.95, 3.5, 0.45, size=22, bold=True, color=ACCENT_TEAL)
    txt(s, b, xi + 0.15, 2.4, 3.5, 0.4, size=13, color=TEXT_GREY)

# AIR result
rect(s, 0.5, 3.2, 5.5, 0.8, BG_CARD)
txt(s, "✅  AI Receptionist 已部署：来电等待 12min → 90sec", 0.65, 3.3, 5.2, 0.55, size=14, color=GREEN_OK)

txt(s, "但路由之后，每天还在浪费：", 0.5, 4.25, 10, 0.45, size=16, color=TEXT_GREY)

# Waste rows
wastes = [
    ("CSR 手动派单 SMS", "5 min × 30 单 × 33 店", "= 82 staff-hours/天"),
    ("Lowe's 传真（完工表单）", "8 min × 31 份 × 33 店", "= 145 staff-hours/天"),
    ("客户投诉首次响应", "平均 30–45 分钟", "≈ 1,000 min/天"),
]
for i, (a, b, c) in enumerate(wastes):
    yi = 4.75 + i * 0.6
    txt(s, a, 0.5, yi, 4.0, 0.5, size=15, color=TEXT_WHITE)
    txt(s, b, 4.5, yi, 3.5, 0.5, size=14, color=TEXT_GREY)
    txt(s, c, 8.2, yi, 4.5, 0.5, size=15, bold=True, color=YELLOW)

add_notes(s, "Keller 是 RC 的真实客户，公开 Case Study 可查。AIR 解决了前台，我们解决后台每一步。")

# ── Slide 4: Insight ─────────────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_BLUE)

txt(s, "洞察", 0.5, 0.3, 3, 0.5, size=13, color=ACCENT_BLUE, bold=True)
txt(s, "三种协作，一个空间", 0.5, 0.65, 10, 0.9, size=42, bold=True)

modes = [
    ("人 ↔ 人", "RC Team Messaging 已有", ACCENT_TEAL, "✅ 现有"),
    ("人 ↔ Agent", "Agent 作为正式团队成员，可被 @", ACCENT_BLUE, "🔨 我们建"),
    ("Agent ↔ Agent", "Agent 之间自动路由任务", GREEN_OK, "🔨 我们建"),
]
for i, (title, desc, color, badge) in enumerate(modes):
    xi = 0.5 + i * 4.2
    rect(s, xi, 1.9, 3.9, 1.6, BG_CARD)
    txt(s, title, xi + 0.15, 2.0, 3.6, 0.55, size=22, bold=True, color=color)
    txt(s, desc, xi + 0.15, 2.55, 3.6, 0.65, size=13, color=TEXT_GREY, wrap=True)
    txt(s, badge, xi + 2.5, 1.92, 1.3, 0.35, size=11, color=color)

divider(s, 3.9)

# Comparison
txt(s, "MuleRun Messages（刚发布）验证了这个方向，但：", 0.5, 4.1, 12, 0.45, size=15, color=TEXT_GREY)

rows = [
    ("MuleRun 的 Agent", "在线程里 讨论", RED_WARN),
    ("AgentRun 的 Agent", "在线程里 执行", GREEN_OK),
]
for i, (who, what, color) in enumerate(rows):
    yi = 4.65 + i * 0.7
    txt(s, who, 0.5, yi, 4.5, 0.55, size=17, color=TEXT_WHITE)
    txt(s, what, 5.0, yi, 5.0, 0.55, size=17, bold=True, color=color)

txt(s, "护城河：RC 的 SMS · Fax · Phone API  —  竞争对手没有",
    0.5, 6.3, 12, 0.5, size=15, bold=True, color=ACCENT_TEAL)

add_notes(s, "MuleRun 刚发布 Messages 验证了市场。但他们没有 RC 的通信 API 深度。这是我们的护城河。")

# ── Slide 5: What We Built ───────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, GREEN_OK)

txt(s, "我们建了什么", 0.5, 0.3, 4, 0.5, size=13, color=GREEN_OK, bold=True)
txt(s, "AgentRun = 两个核心能力", 0.5, 0.65, 12, 0.9, size=38, bold=True)

# Left card
rect(s, 0.5, 1.85, 5.9, 4.5, BG_CARD)
txt(s, "① 角色定义的 Agent", 0.7, 2.0, 5.5, 0.55, size=20, bold=True, color=ACCENT_BLUE)
items_l = [
    "SOUL — 身份 · 规则 · 不可绕过的硬约束",
    "Skills — 独立工作流模块（SKILL.md）",
    "Domain Memory — 业务知识（§ 分隔）",
    "Owner Memory — 个人偏好积累",
    "学习循环 — 从实操创建新技能",
    "",
    "8 个角色 Bot  ·  10 个 Skill 模块",
]
for i, it in enumerate(items_l):
    color = ACCENT_TEAL if it.startswith("8") else TEXT_WHITE
    txt(s, it, 0.75, 2.6 + i * 0.52, 5.5, 0.48, size=13 if it.startswith("8") else 14,
        color=color, bold=it.startswith("8"))

# Right card
rect(s, 7.0, 1.85, 5.9, 4.5, BG_CARD)
txt(s, "② 人→Agent / Agent→Agent", 7.2, 2.0, 5.5, 0.55, size=20, bold=True, color=GREEN_OK)
items_r = [
    "同一 RC 频道，三种协作共存",
    "Sarah 可以直接 @tom-bot",
    "karen-bot 自动路由给 tom-bot",
    "Agent 路由不需要人中转",
    "12 条已定义路由事件",
    "",
    "完整 Agent 协作网络",
]
for i, it in enumerate(items_r):
    color = ACCENT_TEAL if it.startswith("完整") else TEXT_WHITE
    txt(s, it, 7.2, 2.6 + i * 0.52, 5.5, 0.48, size=13 if it.startswith("完整") else 14,
        color=color, bold=it.startswith("完整"))

txt(s, "基础设施：FIJI  ·  AVA Control Plane  ·  K8S  ·  RingClaw  ·  RC APIs",
    0.5, 6.65, 12.3, 0.45, size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)

add_notes(s, "两个能力相互依存：角色 Bot 是节点（深度），Agent-to-Agent 是边（宽度）。缺一不可。")

# ── Slide 6: DEMO 1 – Dispatch ───────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_TEAL)

label(s, "  LIVE DEMO 1  ", 0.5, 0.3, 1.8, 0.38, ACCENT_TEAL, BG_DARK, 11)
txt(s, "30 秒派单闭环", 2.5, 0.25, 8, 0.55, size=34, bold=True)

# Input
rect(s, 0.5, 1.1, 12.3, 1.1, BG_CARD)
txt(s, "Sarah:", 0.7, 1.18, 1.2, 0.38, size=14, color=ACCENT_TEAL, bold=True)
txt(s, "@sarah-bot  dispatch A8821 to Mike, tomorrow 10am,\n"
    "1234 Main St Atlanta GA 30309, Engineered Oak 850sqft, customer Jenkins +1 404-555-0199",
    1.9, 1.15, 10.7, 0.9, size=14, color=TEXT_WHITE)

txt(s, "↓  3 秒", 0.5, 2.35, 2, 0.4, size=16, color=TEXT_DIM, bold=True)

# Output box
rect(s, 0.5, 2.85, 12.3, 1.6, BG_CARD)
txt(s, "sarah-bot:", 0.7, 2.93, 2.2, 0.38, size=14, color=ACCENT_BLUE, bold=True)

results = [
    ("✅  Task #T992 created — Mike Reyes · due 06/04 10:00", GREEN_OK),
    ("✅  SMS → Mike +14045550211 · delivered", GREEN_OK),
    ("⏳  30min 无 CONFIRM → 自动提醒（cron）", YELLOW),
]
for i, (r, c) in enumerate(results):
    txt(s, r, 2.9, 2.95 + i * 0.43, 10.2, 0.4, size=14, color=c, bold=(i < 2))

# Mike phone
rect(s, 0.5, 4.7, 5.8, 1.5, BG_CARD)
txt(s, "📱  Mike 手机收到", 0.7, 4.78, 5.2, 0.4, size=13, color=TEXT_GREY)
txt(s, '"Install #A8821 06/04 10am.\nAddress: 1234 Main St Atlanta GA 30309\nReply CONFIRM to acknowledge."',
    0.7, 5.18, 5.2, 1.0, size=12, color=TEXT_WHITE, italic=True)

# Metrics
rect(s, 7.0, 4.7, 5.8, 1.5, BG_CARD)
txt(s, "之前", 7.2, 4.82, 2.5, 0.4, size=16, bold=True, color=TEXT_GREY)
txt(s, "5 分钟手工操作", 7.2, 5.2, 3.5, 0.45, size=18, color=TEXT_GREY)
txt(s, "现在", 9.8, 4.82, 2.5, 0.4, size=16, bold=True, color=GREEN_OK)
txt(s, "3 秒", 9.8, 5.2, 2.5, 0.45, size=36, bold=True, color=GREEN_OK)

add_notes(s, "演示要点：说一句话，bot 完成：目录查询 Mike 手机 → Task 创建 → SMS 发出 → CONFIRM 跟踪启动。")

# ── Slide 7: DEMO 2 – Complaint Multi-Agent ──────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, RED_WARN)

label(s, "  LIVE DEMO 2  ", 0.5, 0.3, 1.8, 0.38, RED_WARN, BG_DARK, 11)
txt(s, "投诉 → 17 分钟关闭（vs 30-45 分钟）", 2.5, 0.25, 10, 0.55, size=30, bold=True)

timeline = [
    ("10:02", "客户 SMS 到达", '"Crew didn\'t show up. Worst service ever!!!"', TEXT_GREY),
    ("10:03", "sarah-bot 自动", "安抚 SMS 发出 ✅  +  URGENT Task 创建  +  路由 tom-bot", GREEN_OK),
    ("10:05", "tom-bot 自动", "Call Log 查询 + 派工记录对比 → 调查结论发出（无需 Tom 操作）", ACCENT_TEAL),
    ("10:17", "Tom 介入", "读调查结论，致电 Mike 确认（30 秒）", TEXT_WHITE),
    ("10:19", "sarah-bot 执行", "道歉 SMS + $50 credit 发出 ✅", GREEN_OK),
]
for i, (t, who, what, color) in enumerate(timeline):
    yi = 1.1 + i * 1.1
    rect(s, 0.5, yi, 1.0, 0.85, BG_CARD)
    txt(s, t, 0.52, yi + 0.22, 0.96, 0.45, size=13, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    rect(s, 1.65, yi, 2.4, 0.85, BG_CARD)
    txt(s, who, 1.75, yi + 0.2, 2.2, 0.5, size=13, bold=True, color=color)

    rect(s, 4.2, yi, 8.6, 0.85, BG_CARD)
    txt(s, what, 4.35, yi + 0.18, 8.3, 0.55, size=13, color=TEXT_WHITE)

txt(s, "客户自动安抚：51 秒  ·  人只做最终决策  ·  全程审计轨迹", 0.5, 6.75, 12, 0.45,
    size=14, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_notes(s, "关键：10:05 tom-bot 自动调查，Tom 没有操作任何东西。Agent-to-Agent 路由让调查自动启动。")

# ── Slide 8: DEMO 3 – Phone Call ────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, GREEN_OK)

label(s, "  LIVE DEMO 3  ", 0.5, 0.3, 1.8, 0.38, GREEN_OK, BG_DARK, 11)
txt(s, "说\"打过去\"，FIJI 真的拨出电话", 2.5, 0.25, 10, 0.55, size=30, bold=True)

# Conversation
rect(s, 0.5, 1.1, 12.3, 0.75, BG_CARD)
txt(s, "Beth:", 0.7, 1.22, 1.5, 0.4, size=14, color=YELLOW, bold=True)
txt(s, '"Lowe\'s HQ 未接来电，打过去"', 2.2, 1.22, 10, 0.4, size=16, color=TEXT_WHITE)

# Arrow
txt(s, "↓", 6.5, 1.95, 1.5, 0.5, size=24, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Flow
flow_steps = [
    ("beth-bot", "ACTION:PHONE_CALL to=+19195550188", ACCENT_BLUE),
    ("Control Plane", "记录 action_event：{type: PHONE_CALL, status: client_action_required, details: {make_call}}", BG_CARD),
    ("FIJI AvaClientActionBridge", "轮询到 event  →  executeAvaClientAction()", ACCENT_TEAL),
    ("FIJI directCall", "以 Beth 当前登录身份拨出  ✅", GREEN_OK),
]
for i, (step, action, color) in enumerate(flow_steps):
    yi = 2.55 + i * 0.9
    rect(s, 0.5, yi, 2.5, 0.75, color)
    txt(s, step, 0.6, yi + 0.18, 2.3, 0.45, size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    txt(s, action, 3.2, yi + 0.18, 9.8, 0.5, size=13, color=TEXT_WHITE)
    if i < 3:
        txt(s, "↓", 1.5, yi + 0.77, 1.0, 0.3, size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

divider(s, 6.35, ACCENT_TEAL)

# Competitor comparison
cols = [("Copilot", "✗", RED_WARN), ("Gemini", "✗", RED_WARN), ("AgentRun", "✅", GREEN_OK)]
for i, (name, mark, color) in enumerate(cols):
    xi = 1.5 + i * 3.8
    txt(s, name, xi, 6.5, 3.0, 0.45, size=16, color=TEXT_GREY, align=PP_ALIGN.CENTER)
    txt(s, mark, xi, 6.92, 3.0, 0.4, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)

txt(s, "护城河：RC Phone API  —  竞争对手没有", 0.5, 7.05, 12, 0.35,
    size=13, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)

add_notes(s, "这是最有力的 demo 时刻。Beth 说一句话，FIJI 真的打出电话。这是 RC 独一无二的能力。")

# ── Slide 9: Role Bot Value ──────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_BLUE)

txt(s, "角色定义 Bot 的独特价值", 0.5, 0.3, 12, 0.7, size=36, bold=True)

# Comparison table
headers = ["维度", "通用 AI（Copilot / GPT）", "AgentRun 角色 Bot"]
widths  = [2.8, 4.2, 4.2]
xs      = [0.5, 3.4, 7.7]

for j, (h, w, x) in enumerate(zip(headers, widths, xs)):
    color = ACCENT_BLUE if j == 2 else (TEXT_GREY if j == 1 else TEXT_WHITE)
    rect(s, x, 1.2, w, 0.55, BG_CARD if j > 0 else BG_DARK)
    txt(s, h, x + 0.1, 1.28, w - 0.2, 0.42, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

rows_data = [
    ("业务上下文",       "每次对话从零解释",        "SOUL + DOMAIN.md 内置"),
    ("规则执行",         "可被提示词绕过",           "SOUL 硬规则，真正不可绕"),
    ("组织结构",         "不知道找谁升级",           "知道出了事该找谁"),
    ("员工离职",         "知识随人消失",             "机构记忆沉淀在 Bot 里"),
    ("持续学习",         "每次重新开始",             "Skills 学习循环，越用越准"),
]
for i, (dim, generic, ours) in enumerate(rows_data):
    yi = 1.85 + i * 0.85
    bg_row = BG_CARD if i % 2 == 0 else BG_DARK
    for j, (cell, x, w, color) in enumerate([
        (dim, 0.5, 2.8, TEXT_GREY),
        (generic, 3.4, 4.2, TEXT_GREY),
        (ours, 7.7, 4.2, GREEN_OK),
    ]):
        rect(s, x, yi, w, 0.72, bg_row)
        txt(s, cell, x + 0.12, yi + 0.18, w - 0.2, 0.42, size=13, color=color)

txt(s, "HR bot 永远不会把员工请假原因发到群里  ·  Finance bot 永远不会在没有审批的情况下付款",
    0.5, 6.8, 12.3, 0.45, size=13, color=ACCENT_TEAL, italic=True, align=PP_ALIGN.CENTER)

add_notes(s, "角色 Bot 和通用 AI 的根本差异：不是模型更好，是它真正知道你是谁、你的业务规则、你的团队结构。")

# ── Slide 10: Business Value ─────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, YELLOW)

txt(s, "商业价值", 0.5, 0.3, 3, 0.5, size=13, color=YELLOW, bold=True)
txt(s, "Keller · 33 门店 · 150 名 CSR", 0.5, 0.65, 12, 0.8, size=36, bold=True)

metrics = [
    ("一张派单",         "5 min",     "3 秒",       "99%↓"),
    ("Lowe's 传真批次",  "4-5 小时",  "等 49 分钟", "83%↓"),
    ("投诉首次响应",     "30-45 min", "51 秒",      "97%↓"),
    ("每日门店摘要×33",  "16.5h/天",  "0",          "100%↓"),
    ("员工请假申请",     "邮件往返",  "当天 DM",    "当天完成"),
    ("分包商月度付款",   "半天 Excel","Bot 汇总",   "90%↓"),
]
for i, (scene, before, after, delta) in enumerate(metrics):
    xi = 0.5 + (i % 2) * 6.4
    yi = 1.75 + (i // 2) * 1.6
    rect(s, xi, yi, 6.0, 1.35, BG_CARD)
    txt(s, scene, xi + 0.18, yi + 0.12, 5.6, 0.45, size=15, bold=True, color=TEXT_WHITE)
    txt(s, f"之前：{before}", xi + 0.18, yi + 0.58, 2.8, 0.4, size=13, color=TEXT_GREY)
    txt(s, f"之后：{after}", xi + 3.0, yi + 0.58, 2.8, 0.4, size=13, color=GREEN_OK, bold=True)
    txt(s, delta, xi + 4.5, yi + 0.12, 1.4, 0.4, size=14, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

txt(s, "≈  80+ staff-hours/天 节省  ·  相当于每天少雇 10 个人做重复性工作",
    0.5, 7.02, 12.3, 0.4, size=14, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_notes(s, "数字基于 Keller 规模估算。最有力的是 Lowe's 传真：之前需要人一份一份手发，现在 Karen 输入一行命令等 49 分钟。")

# ── Slide 11: Architecture ───────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, TEXT_DIM)

txt(s, "技术架构", 0.5, 0.3, 4, 0.5, size=13, color=TEXT_GREY, bold=True)
txt(s, "基于现有 RC 平台，增量构建", 0.5, 0.65, 12, 0.8, size=34, bold=True)

layers = [
    ("FIJI UI",          "AvaBotOnboarding  ·  AvaClientActionBridge → PHONE_CALL  ✅",     ACCENT_TEAL),
    ("AVA Control Plane","Bot Registry  ·  Token Pool  ·  Lifecycle  ·  Action Events  ✅", ACCENT_BLUE),
    ("RingClaw Runtime", "SOUL + SKILL  ·  Agent-to-Agent  ·  RC Actions  ·  K8S Pod/Bot",  ACCENT_BLUE),
    ("RC Platform APIs", "SMS  ·  Fax  ·  Phone  ·  Task  ·  Note  ·  Video  ·  CallLog  ✅", GREEN_OK),
]
for i, (layer, detail, color) in enumerate(layers):
    yi = 1.8 + i * 1.25
    rect(s, 0.5, yi, 1.8, 1.0, color)
    txt(s, layer, 0.6, yi + 0.25, 1.6, 0.55,
        size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    rect(s, 2.45, yi, 10.35, 1.0, BG_CARD)
    txt(s, detail, 2.6, yi + 0.28, 10.1, 0.5, size=14, color=TEXT_WHITE)
    if i < 3:
        txt(s, "↕", 1.2, yi + 1.02, 0.7, 0.35, size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

txt(s, "新增代码量：~500 行  ·  其余全部复用现有平台能力",
    0.5, 7.08, 12.3, 0.35, size=13, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)

add_notes(s, "评委关注：这不是从零搭建，是在现有 RC 平台上增量构建。PHONE_CALL bridge 已经 merge 到代码里了。")

# ── Slide 12: What's Done ────────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, GREEN_OK)

txt(s, "Hackathon 完成情况", 0.5, 0.3, 5, 0.5, size=13, color=GREEN_OK, bold=True)
txt(s, "已完成  vs  剩余", 0.5, 0.65, 12, 0.8, size=36, bold=True)

done = [
    "8 个角色 Bot 完整 SOUL + config 设计",
    "10 个 Skill 模块（SKILL.md，Hermes 格式）",
    "12 条 Agent-to-Agent 路由事件定义",
    "PHONE_CALL → FIJI directCall 完整打通  ✅",
    "AVA Control Plane 全量 API（已实现）",
    "Keller POC 8 周实施路径",
    "编译错误修复（代码可构建）",
]
todo = [
    "Inbound SMS/Fax wire（~150 行）",
    "OOB 审批反向通道（~200 行）",
    "Frozen Snapshot 优化（~100 行）",
]

rect(s, 0.5, 1.75, 7.5, 5.0, BG_CARD)
txt(s, "✅  已完成", 0.7, 1.85, 7.0, 0.45, size=16, bold=True, color=GREEN_OK)
for i, item in enumerate(done):
    txt(s, f"  ✓  {item}", 0.7, 2.42 + i * 0.6, 7.1, 0.52, size=13, color=TEXT_WHITE)

rect(s, 8.3, 1.75, 4.7, 5.0, BG_CARD)
txt(s, "📌  剩余（代码层）", 8.5, 1.85, 4.4, 0.45, size=16, bold=True, color=YELLOW)
for i, item in enumerate(todo):
    txt(s, f"  ·  {item}", 8.5, 2.42 + i * 0.8, 4.3, 0.7, size=13, color=TEXT_WHITE, wrap=True)

txt(s, "剩余约 500 行代码  ·  技术路径清晰，无未知风险",
    0.5, 7.08, 12.3, 0.35, size=13, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_notes(s, "评委问'你们真正建了什么'时：Control Plane 是生产代码，PHONE_CALL bridge 已 merge，Keller 设计完整可实施。")

# ── Slide 13: Why Now / Why RC ───────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_TEAL)

txt(s, "为什么是现在，为什么是 RC", 0.5, 0.3, 12, 0.7, size=36, bold=True)

reasons = [
    (
        "1",
        "MuleRun 验证了市场",
        "刚发布 Messages，企业 AI 协作需求已被验证\n但 MuleRun 没有 RC 的通信 API",
        ACCENT_BLUE,
    ),
    (
        "2",
        "RC 是唯一有这些 API 的平台",
        "SMS · Fax · Phone · Task · Note · Video\n竞争对手做「讨论」，我们做「执行」",
        ACCENT_TEAL,
    ),
    (
        "3",
        "AI Receptionist 已教育了市场",
        "Keller：12min → 90sec，客户已接受 AI\n下一步：内部协作也要 AI",
        GREEN_OK,
    ),
]
for i, (num, title, desc, color) in enumerate(reasons):
    yi = 1.4 + i * 1.8
    rect(s, 0.5, yi, 0.9, 1.4, color)
    txt(s, num, 0.52, yi + 0.38, 0.86, 0.7, size=32, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.55, yi, 11.2, 1.4, BG_CARD)
    txt(s, title, 1.75, yi + 0.18, 10.8, 0.55, size=20, bold=True, color=color)
    txt(s, desc, 1.75, yi + 0.72, 10.8, 0.58, size=14, color=TEXT_GREY, wrap=True)

divider(s, 7.0, ACCENT_TEAL)
txt(s, "RingEX  +  AI Receptionist（前台）  +  AgentRun（后台）  =  企业通信完整 AI 栈",
    0.5, 7.1, 12.3, 0.35, size=14, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

add_notes(s, "这三个时机叠加才有今天的窗口。RC 的护城河是真实的，不是 story。")

# ── Slide 14: Roadmap ────────────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_BLUE)

txt(s, "路线图", 0.5, 0.3, 3, 0.5, size=13, color=ACCENT_BLUE, bold=True)
txt(s, "从 POC 到产品", 0.5, 0.65, 12, 0.8, size=36, bold=True)

phases = [
    ("Q3 2026\nPOC",
     ["Keller 5 人 Atlanta 试点",
      "Case 1-6 全部跑通",
      "Inbound wire → Case 7/8 上线",
      "Keller 扩展 33 门店"],
     ACCENT_BLUE),
    ("Q4 2026\n深度",
     ["Frozen Snapshot + Skills 学习循环",
      "OOB 审批入 FIJI（不需要 terminal）",
      "HR/Finance Bot 全量上线"],
     ACCENT_TEAL),
    ("2027 Q1\n产品化",
     ["SOUL/Memory 管理 UI",
      "Skills Marketplace",
      "其他 RC 客户自助注册",
      "App Gallery 上架"],
     GREEN_OK),
]
for i, (phase, items, color) in enumerate(phases):
    xi = 0.5 + i * 4.3
    rect(s, xi, 1.75, 4.0, 0.75, color)
    txt(s, phase, xi + 0.2, 1.85, 3.6, 0.6, size=16, bold=True, color=TEXT_WHITE)
    rect(s, xi, 2.6, 4.0, 3.8, BG_CARD)
    for j, item in enumerate(items):
        txt(s, f"·  {item}", xi + 0.18, 2.75 + j * 0.82, 3.7, 0.72, size=14, color=TEXT_WHITE, wrap=True)

divider(s, 6.85)
txt(s, "North Star：每个 RC 企业用户，都有一个懂自己业务的 Agent 同事",
    0.5, 7.0, 12.3, 0.42, size=14, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_notes(s, "路线图是实事求是的。POC 8 周完成后才谈扩展，不是愿景 PPT。")

# ── Slide 15: Closing ────────────────────────────────────────────

s = add_slide()
bg(s)
rect(s, 0, 0, 0.12, 7.5, ACCENT_BLUE)

# Contrast statement
txt(s, "其他 AI 助手帮你", 0.5, 1.4, 12.3, 0.85,
    size=44, color=TEXT_GREY, align=PP_ALIGN.CENTER)
txt(s, "想", 0.5, 2.22, 12.3, 1.2,
    size=96, bold=True, color=TEXT_GREY, align=PP_ALIGN.CENTER)

divider(s, 3.9, ACCENT_BLUE, 12.5)

txt(s, "AgentRun 帮你", 0.5, 4.0, 12.3, 0.85,
    size=44, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txt(s, "做", 0.5, 4.8, 12.3, 1.2,
    size=96, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

txt(s, "因为我们有 RingCentral  ·  SMS 真发出去  ·  传真真送到  ·  电话真打出来",
    0.5, 6.5, 12.3, 0.5, size=16, color=TEXT_GREY, align=PP_ALIGN.CENTER)

txt(s, "Agent 不只是你的助手，是你团队里的同事",
    0.5, 7.05, 12.3, 0.35, size=14, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_notes(s, "结尾要干脆有力。停顿在「做」字上一秒，然后说：因为我们有 RingCentral。")

# ── Save ──────────────────────────────────────────────────────────

out = "/Users/summer.gan/Work/code/Aplication/ringclaw/outputs/keller-poc/AgentRun-Hackathon-Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
